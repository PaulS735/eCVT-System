"""
Add new slides to the Critical Design Review presentation to address reviewer feedback.
Inserts 8 new slides at specified locations without deleting existing slides.
Inserts are done in order from lowest to highest original index to keep offset tracking correct.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from copy import deepcopy
from lxml import etree
import os

ORIGINAL_PATH = r"C:\Users\scpks\Downloads\Critical Design Review.pptx"
OUTPUT_PATH = r"C:\Users\scpks\Documents\Data Logging Teensy\Critical Design Review.pptx"

prs = Presentation(ORIGINAL_PATH)

slide_count = len(prs.slides)
if slide_count != 44:
    raise RuntimeError(f"Expected 44 original slides but found {slide_count}.")

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

LAYOUT = prs.slide_layouts[1]  # "Title and Content"

# Styling constants
TITLE_FONT_SIZE = Pt(28)
BODY_FONT_SIZE = Pt(14)
SMALL_FONT_SIZE = Pt(12)
TABLE_FONT_SIZE = Pt(11)
TABLE_HEADER_FONT_SIZE = Pt(12)
NOTE_FONT_SIZE = Pt(11)
MONO_FONT = "Consolas"
BODY_FONT = "Calibri"

TITLE_LEFT = 700635
TITLE_TOP = 914400
TITLE_WIDTH = 10691265
TITLE_HEIGHT = 1307592

CONTENT_LEFT = 700635
CONTENT_TOP = 2221992
CONTENT_WIDTH = 10691265
CONTENT_HEIGHT = 3739896


def insert_slide_at(prs, index, layout):
    """Insert a new slide at the given index (0-based) in the presentation."""
    slide = prs.slides.add_slide(layout)
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    new_slide_elem = slides_list[-1]
    xml_slides.remove(new_slide_elem)
    if index >= len(slides_list) - 1:
        xml_slides.append(new_slide_elem)
    else:
        xml_slides.insert(index, new_slide_elem)
    return slide


def set_title(slide, title_text):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            shape.text = title_text
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = TITLE_FONT_SIZE
            return shape
    return None


def clear_body_placeholder(slide):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            shape.text_frame.clear()
            return shape
    return None


def add_textbox(slide, left, top, width, height):
    from pptx.util import Emu
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    return txBox


def add_body_paragraph(tf, text, font_size=BODY_FONT_SIZE, bold=False, italic=False,
                       font_name=None, bullet=False, space_before=Pt(2), space_after=Pt(2),
                       color=None, level=0):
    p = tf.add_paragraph()
    p.space_before = space_before
    p.space_after = space_after
    p.level = level
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    if font_name:
        run.font.name = font_name
    else:
        run.font.name = BODY_FONT
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return p


def add_table(slide, rows, cols, left, top, width, height):
    table_shape = slide.shapes.add_table(rows, cols, Emu(left), Emu(top), Emu(width), Emu(height))
    return table_shape.table


def set_cell(table, row, col, text, font_size=TABLE_FONT_SIZE, bold=False, alignment=PP_ALIGN.LEFT):
    cell = table.cell(row, col)
    cell.text = text
    for para in cell.text_frame.paragraphs:
        para.alignment = alignment
        for run in para.runs:
            run.font.size = font_size
            run.font.bold = bold
            run.font.name = BODY_FONT


def shade_row(table, row_idx, color):
    for col_idx in range(len(table.columns)):
        cell = table.cell(row_idx, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = color


# ============================================================
# Insert slides from LOWEST original index to HIGHEST
# Original slide numbers (1-based):
#   After 10 -> Actuator Speed  (insert at index 10)
#   After 13 -> Power Budget    (insert at index 13)
#   After 15 -> Electrical Schematic (insert at index 15)
#   After 21 -> State Machine   (insert at index 21)
#   After state machine -> Control Loop Detail (insert at index 22)
#   After 32 -> FoS Methodology (insert at index 32)
#   After 35 -> FMEA Summary   (insert at index 35)
#   Before Appendix (44) -> Budget (insert at index 43)
# ============================================================

offset = 0

# ============================================================
# 1. Actuator Speed Verification (insert after original slide 10)
# ============================================================
insert_idx = 10 + offset  # 10
slide = insert_slide_at(prs, insert_idx, LAYOUT)
set_title(slide, "Key Design Decisions: Actuator Speed Verification")

body = clear_body_placeholder(slide)
tf = body.text_frame
tf.word_wrap = True

# Specs table
tbl = add_table(slide, 5, 3,
                left=CONTENT_LEFT, top=CONTENT_TOP,
                width=int(CONTENT_WIDTH * 0.55), height=Emu(900000))

specs_headers = ["Parameter", "Value", "Requirement"]
for i, h in enumerate(specs_headers):
    set_cell(tbl, 0, i, h, font_size=TABLE_HEADER_FONT_SIZE, bold=True)
shade_row(tbl, 0, RGBColor(0x1F, 0x49, 0x7D))
for ci in range(3):
    for para in tbl.cell(0, ci).text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

specs_data = [
    ["Stroke", "100mm", "--"],
    ["No-load speed", "250ms full stroke", "--"],
    ["Loaded speed", "~500ms full stroke", "--"],
    ["Typical shift (25mm)", "125ms under load", "SR_07: 250ms"],
]
for r, row_data in enumerate(specs_data):
    for c, val in enumerate(row_data):
        set_cell(tbl, r + 1, c, val)
    if r % 2 == 0:
        shade_row(tbl, r + 1, RGBColor(0xE8, 0xEE, 0xF4))

# Compliance summary next to table
comp_box = add_textbox(slide, CONTENT_LEFT + int(CONTENT_WIDTH * 0.57),
                       CONTENT_TOP, int(CONTENT_WIDTH * 0.43), Emu(900000))
ctf = comp_box.text_frame
ctf.word_wrap = True
p = ctf.paragraphs[0]
run = p.add_run()
run.text = "Compliance Summary"
run.font.size = Pt(14)
run.font.bold = True
run.font.name = BODY_FONT

add_body_paragraph(ctf, "SR_07 (250ms shift): PASS", font_size=Pt(13), bold=True,
                   color=RGBColor(0x00, 0x70, 0x30))
add_body_paragraph(ctf, "  25mm step = 125ms (50% margin)", font_size=Pt(12))
add_body_paragraph(ctf, "SR_16 (100ms command): PASS", font_size=Pt(13), bold=True,
                   color=RGBColor(0x00, 0x70, 0x30))
add_body_paragraph(ctf, "  Command cycle < 15ms", font_size=Pt(12))

# Voltage divider and ADC analysis below
detail_box = add_textbox(slide, CONTENT_LEFT, CONTENT_TOP + 1050000,
                         CONTENT_WIDTH, Emu(2600000))
dtf = detail_box.text_frame
dtf.word_wrap = True

add_body_paragraph(dtf, "Voltage Divider & ADC Analysis", font_size=Pt(15), bold=True,
                   space_before=Pt(0))
add_body_paragraph(dtf, "Actuator feedback: 0-5V output through 10k/8.2k voltage divider", font_size=Pt(13))
add_body_paragraph(dtf, "  Vout = 5V \u00d7 8.2k / (10k + 8.2k) = 2.25V (1.05V below 3.3V ADC ref = safe)", font_size=Pt(12),
                   font_name=MONO_FONT)
add_body_paragraph(dtf, "  ADC range: 0-2794 counts (of 4095 max on 12-bit ADC)", font_size=Pt(12),
                   font_name=MONO_FONT)
add_body_paragraph(dtf, "  Resolution: 100mm / 2794 = 0.036 mm/count (28 counts/mm)", font_size=Pt(12),
                   font_name=MONO_FONT)
add_body_paragraph(dtf, "", font_size=Pt(4))
add_body_paragraph(dtf, "Deadband: 50 ADC counts = 1.79mm physical", font_size=Pt(13), bold=True)
add_body_paragraph(dtf, "  Exceeds mechanical precision needs of CVT sheave system (belt stretch, backlash >> 1.79mm)",
                   font_size=Pt(12))

offset += 1

# ============================================================
# 2. Power Budget (insert after original slide 13)
# ============================================================
insert_idx = 13 + offset  # 14
slide = insert_slide_at(prs, insert_idx, LAYOUT)
set_title(slide, "Electrical Interfaces: Power Budget Analysis")

body = clear_body_placeholder(slide)
tf = body.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = ""

tbl = add_table(slide, 10, 5,
                left=CONTENT_LEFT, top=CONTENT_TOP,
                width=CONTENT_WIDTH, height=Emu(1800000))

headers = ["Component", "Supply", "Typical (A)", "Peak (A)", "Notes"]
for i, h in enumerate(headers):
    set_cell(tbl, 0, i, h, font_size=TABLE_HEADER_FONT_SIZE, bold=True)
shade_row(tbl, 0, RGBColor(0x1F, 0x49, 0x7D))
for ci in range(5):
    for para in tbl.cell(0, ci).text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

data = [
    ["PQ12-100-6-R Actuator", "12V direct", "2.00", "10-15 (stall)", "Relay H-bridge driven"],
    ["Raspberry Pi 4B", "5V (buck)", "1.00", "3.00", "USB-C powered"],
    ["Teensy 4.1", "5V (Pi USB)", "0.10", "0.10", "Includes sensor draw"],
    ["Relay Module", "5V (buck)", "0.05", "0.10", "2 coils + optocouplers"],
    ["Hall Effect Sensor", "3.3V (Teensy)", "0.01", "0.01", "A3144 equivalent"],
    ["Potentiometer", "3.3V (Teensy)", "0.001", "0.001", "10k divider"],
    ["LoRa RFM95W (future)", "3.3V (reg)", "0.012", "0.12", "Not yet implemented"],
    ["TOTAL at 12V input", "", "2.57", "16.61", "Referred through buck"],
]
for r, row_data in enumerate(data):
    for c, val in enumerate(row_data):
        bold = (r == len(data) - 1)
        set_cell(tbl, r + 1, c, val, bold=bold)
    if r % 2 == 0:
        shade_row(tbl, r + 1, RGBColor(0xE8, 0xEE, 0xF4))
shade_row(tbl, 9, RGBColor(0xD6, 0xE4, 0xF0))

col_widths = [0.28, 0.17, 0.15, 0.17, 0.23]
for i, pct in enumerate(col_widths):
    tbl.columns[i].width = int(CONTENT_WIDTH * pct)

summary_box = add_textbox(slide, CONTENT_LEFT, CONTENT_TOP + 1900000, CONTENT_WIDTH, 1600000)
stf = summary_box.text_frame
stf.word_wrap = True

add_body_paragraph(stf, "Margin Analysis:", font_size=Pt(14), bold=True, space_before=Pt(0))
add_body_paragraph(stf, "Typical: 2.57A of 18A limit = 85.7% headroom (PASS)", font_size=Pt(13))
add_body_paragraph(stf, "Peak (stall): 16.61A of 18A limit = 7.7% margin, transient only (MARGINAL)", font_size=Pt(13))
add_body_paragraph(stf, "Buck converter: 1.15A / 3.0A = 38.3% of rated capacity", font_size=Pt(13))
add_body_paragraph(stf, "Thermal: 1.3W continuous dissipation, viable for 4-hour endurance, no active cooling", font_size=Pt(13))
add_body_paragraph(stf, "Note: Stall current (15A) is transient (<100ms); 20A fuse survives per I\u00b2t curve",
                   font_size=NOTE_FONT_SIZE, italic=True, color=RGBColor(0x66, 0x66, 0x66))

offset += 1

# ============================================================
# 3. Electrical Schematic (insert after original slide 15)
# ============================================================
insert_idx = 15 + offset  # 17
slide = insert_slide_at(prs, insert_idx, LAYOUT)
set_title(slide, "Electrical Interfaces: System Schematic")

body = clear_body_placeholder(slide)
sp = body._element
sp.getparent().remove(sp)

schema_box = add_textbox(slide, CONTENT_LEFT, CONTENT_TOP - 200000, CONTENT_WIDTH, 4200000)
stf = schema_box.text_frame
stf.word_wrap = True

schematic_lines = [
    "ALTERNATOR 12V",
    "    |",
    "    +-- [20A Blade Fuse] --+",
    "    |                      |",
    "    |                  [Buck Converter 12V\u21925V, 3A, 85% eff.]",
    "    |                      |",
    "    |                      +-- 5V USB-C --> Raspberry Pi 4B",
    "    |                      |                    |",
    "    |                      |                    +-- USB-A 5V+Data --> Teensy 4.1",
    "    |                      |                         |",
    "    |                      |                         +-- 3.3V LDO --> Hall Sensor (Pin 2, PULLUP, FALLING ISR)",
    "    |                      |                         +-- 3.3V LDO --> Potentiometer 10k\u03a9 (Pin A8)",
    "    |                      |                         +-- Pin 3 --> Relay IN1 (FWD, active LOW)",
    "    |                      |                         +-- Pin 4 --> Relay IN2 (REV, active LOW)",
    "    |                      |                         +-- Pin A1 <-- Actuator feedback",
    "    |                      |                              (10k/8.2k divider: 5V\u21922.25V)",
    "    |                      |",
    "    |                      +-- 5V VCC --> Relay Module Logic",
    "    |",
    "    +-- [Relay H-Bridge COM] --> PQ12-100-6-R Actuator",
    "         Wires: Red(+12V) / Black(GND) / White(Feedback)",
]

p = stf.paragraphs[0]
run = p.add_run()
run.text = schematic_lines[0]
run.font.size = Pt(10)
run.font.name = MONO_FONT

for line in schematic_lines[1:]:
    add_body_paragraph(stf, line, font_size=Pt(10), font_name=MONO_FONT,
                       space_before=Pt(0), space_after=Pt(0))

add_body_paragraph(stf, "", font_size=Pt(4), space_before=Pt(4), space_after=Pt(0))
add_body_paragraph(stf, "Note: Full harnessing plan with connector types, wire gauges, and routing to be completed during build phase.",
                   font_size=NOTE_FONT_SIZE, italic=True, color=RGBColor(0x66, 0x66, 0x66))

offset += 1

# ============================================================
# 4. State Machine (insert after original slide 21 = Telemetry Data)
# ============================================================
insert_idx = 21 + offset  # 24
slide = insert_slide_at(prs, insert_idx, LAYOUT)
set_title(slide, "Software Architecture: Firmware State Machine")

body = clear_body_placeholder(slide)
sp = body._element
sp.getparent().remove(sp)

state_box = add_textbox(slide, CONTENT_LEFT, CONTENT_TOP - 300000,
                        int(CONTENT_WIDTH * 0.55), 4200000)
stf = state_box.text_frame
stf.word_wrap = False

state_lines = [
    "    +--------+",
    "    |  INIT  |  Configure pins, start serial,",
    "    |        |  attach ISR, relays OFF",
    "    +---+----+",
    "        | (on boot)",
    "        v",
    "    +---+----+",
    "    |  IDLE  | <----+  Waiting for hall pulse",
    "    +---+----+      |",
    "        |           | RPM timeout (1s)",
    "        | Hall      |",
    "        | pulse     |",
    "        v           |",
    "    +---+----+      |",
    "    | RUNNING|------+",
    "    |        |---+",
    "    +--------+   | Fault detected",
    "                  v",
    "    +--------+   +--------+",
    "    |FAIL_   |<--| FAULT  |",
    "    |SAFE    |   |        |",
    "    +--------+   +---+----+",
    "    Power cycle       | Non-critical",
    "    to reset           --> returns to RUNNING",
]

p = stf.paragraphs[0]
run = p.add_run()
run.text = state_lines[0]
run.font.size = Pt(10)
run.font.name = MONO_FONT

for line in state_lines[1:]:
    add_body_paragraph(stf, line, font_size=Pt(10), font_name=MONO_FONT,
                       space_before=Pt(0), space_after=Pt(0))

desc_box = add_textbox(slide, CONTENT_LEFT + int(CONTENT_WIDTH * 0.50),
                       CONTENT_TOP - 300000, int(CONTENT_WIDTH * 0.50), 4200000)
dtf = desc_box.text_frame
dtf.word_wrap = True

p = dtf.paragraphs[0]
run = p.add_run()
run.text = "State Descriptions"
run.font.size = Pt(14)
run.font.bold = True
run.font.name = BODY_FONT

states = [
    ("INIT", "Hardware setup: GPIO config, serial start, ISR attach, relays forced OFF. Runs once on boot."),
    ("IDLE", "Engine not running. Actuator holds/retracts. Waits for first valid hall pulse to enter RUNNING."),
    ("RUNNING", "Active control: calculates RPM, reads preset selector, drives actuator via bang-bang. Telemetry every 50ms."),
    ("FAULT", "Non-critical: log warning, continue degraded. Critical: stop actuator, latch fault code every 1s."),
    ("FAIL_SAFE", "Terminal state. Both relays OFF (actuator coasts). No software exit. Power cycle to reset."),
]

for name, desc in states:
    add_body_paragraph(dtf, name, font_size=Pt(12), bold=True,
                       color=RGBColor(0x1F, 0x49, 0x7D), space_before=Pt(6))
    add_body_paragraph(dtf, desc, font_size=Pt(11), space_before=Pt(0))

add_body_paragraph(dtf, "", font_size=Pt(2))
add_body_paragraph(dtf, "State machine replaces code screenshots for clarity \u2014 full source available in repository.",
                   font_size=NOTE_FONT_SIZE, italic=True, color=RGBColor(0x66, 0x66, 0x66))

offset += 1

# ============================================================
# 5. Control Loop Detail (insert right after state machine = after original slide 22 position)
# We want it right after the state machine slide we just inserted.
# The state machine is at index 24 (0-based), so we insert at 25.
# But we need to use original numbering. The state machine was inserted after orig 21.
# Control loop should go after state machine = after orig 21 + 1 conceptually.
# With offset=4 now, we insert at 22 + offset = 26? No.
# Actually: state machine was inserted at index 24. We want the next one right after it.
# Since state machine is at position 24 (0-based), the next insert should be at 25.
# But our offset is already 4. Let's use the absolute: 21 + offset (which = 25).
# ============================================================
insert_idx = 21 + offset  # 25 - right after the state machine at 24
slide = insert_slide_at(prs, insert_idx, LAYOUT)
set_title(slide, "Software Architecture: Control Loop Detail")

body = clear_body_placeholder(slide)
tf = body.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = ""

add_body_paragraph(tf, "Hall Effect ISR (RPM Measurement)", font_size=Pt(15), bold=True,
                   space_before=Pt(0), color=RGBColor(0x1F, 0x49, 0x7D))
add_body_paragraph(tf, "FALLING edge interrupt on Pin 2; ISR captures pulse delta via micros()",
                   font_size=Pt(13))
add_body_paragraph(tf, "Main loop copies volatile data with interrupts disabled (no torn reads)",
                   font_size=Pt(13))
add_body_paragraph(tf, "RPM = 60,000,000 / (pulseDelta_us \u00d7 magnetsPerRev)",
                   font_size=Pt(13), font_name=MONO_FONT, bold=True)
add_body_paragraph(tf, "1-second timeout with no pulse \u2192 transition to IDLE",
                   font_size=Pt(13))

add_body_paragraph(tf, "", font_size=Pt(4))
add_body_paragraph(tf, "Preset Selection (12-bit ADC on Pin A8)", font_size=Pt(15), bold=True,
                   color=RGBColor(0x1F, 0x49, 0x7D))
add_body_paragraph(tf, "ADC 0-1364: Economy  |  1365-2729: Sport  |  2730-4095: Aggressive",
                   font_size=Pt(13), font_name=MONO_FONT)
add_body_paragraph(tf, "Each preset defines a piecewise linear RPM-to-position curve",
                   font_size=Pt(13))

add_body_paragraph(tf, "", font_size=Pt(4))
add_body_paragraph(tf, "Bang-Bang Actuator Control", font_size=Pt(15), bold=True,
                   color=RGBColor(0x1F, 0x49, 0x7D))
add_body_paragraph(tf, "Error > +50 counts: Drive EXTEND  |  Error < -50 counts: Drive RETRACT  |  Within deadband: STOP",
                   font_size=Pt(13), font_name=MONO_FONT)
add_body_paragraph(tf, "75ms relay deadtime between direction reversals (shoot-through prevention)",
                   font_size=Pt(13))

add_body_paragraph(tf, "", font_size=Pt(4))
add_body_paragraph(tf, "Why Bang-Bang Over PID", font_size=Pt(15), bold=True,
                   color=RGBColor(0x1F, 0x49, 0x7D))

tbl = add_table(slide, 3, 3,
                left=CONTENT_LEFT, top=CONTENT_TOP + 2800000,
                width=CONTENT_WIDTH, height=Emu(700000))

tbl_headers = ["Property", "Bang-Bang + Deadband", "PID + Relay"]
for i, h in enumerate(tbl_headers):
    set_cell(tbl, 0, i, h, font_size=TABLE_HEADER_FONT_SIZE, bold=True)
shade_row(tbl, 0, RGBColor(0x1F, 0x49, 0x7D))
for ci in range(3):
    for para in tbl.cell(0, ci).text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

tbl_data = [
    ["Precision / Tuning", "1.79mm deadband, zero tuning", "Relay quantization limits; 3 gains needed"],
    ["Oscillation / Wear", "None (deadband prevents chatter)", "High risk if mistuned; relay PWM wear"],
]
for r, row_data in enumerate(tbl_data):
    for c, val in enumerate(row_data):
        set_cell(tbl, r + 1, c, val)
    if r % 2 == 0:
        shade_row(tbl, r + 1, RGBColor(0xE8, 0xEE, 0xF4))

offset += 1

# ============================================================
# 6. FoS Methodology (insert after original slide 32)
# ============================================================
insert_idx = 32 + offset  # 37
slide = insert_slide_at(prs, insert_idx, LAYOUT)
set_title(slide, "Analysis: Factor of Safety Methodology")

body = clear_body_placeholder(slide)
tf = body.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = ""

add_body_paragraph(tf, "Design Safety Factors", font_size=Pt(15), bold=True,
                   space_before=Pt(0), color=RGBColor(0x1F, 0x49, 0x7D))

tbl = add_table(slide, 3, 3,
                left=CONTENT_LEFT, top=CONTENT_TOP + 400000,
                width=int(CONTENT_WIDTH * 0.6), height=Emu(550000))

fos_headers = ["Parameter", "Value", "Application"]
for i, h in enumerate(fos_headers):
    set_cell(tbl, 0, i, h, font_size=TABLE_HEADER_FONT_SIZE, bold=True)
shade_row(tbl, 0, RGBColor(0x1F, 0x49, 0x7D))
for ci in range(3):
    for para in tbl.cell(0, ci).text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

fos_data = [
    ["Yield FoS", "1.25", "Increased if load confidence is low"],
    ["Ultimate FoS", "1.5", "All structural components"],
]
for r, row_data in enumerate(fos_data):
    for c, val in enumerate(row_data):
        set_cell(tbl, r + 1, c, val, bold=(c == 1))
    shade_row(tbl, r + 1, RGBColor(0xE8, 0xEE, 0xF4))

formula_box = add_textbox(slide, CONTENT_LEFT, CONTENT_TOP + 1100000,
                          CONTENT_WIDTH, 1200000)
ftf = formula_box.text_frame
ftf.word_wrap = True

add_body_paragraph(ftf, "Margin Calculation", font_size=Pt(15), bold=True,
                   space_before=Pt(0), color=RGBColor(0x1F, 0x49, 0x7D))
add_body_paragraph(ftf, "Margin to Ultimate = (Ultimate Strength / (FoS \u00d7 Applied Load)) - 1",
                   font_size=Pt(14), font_name=MONO_FONT, bold=True)
add_body_paragraph(ftf, "Margin must be > 0 for all structural components",
                   font_size=Pt(13))
add_body_paragraph(ftf, "", font_size=Pt(4))
add_body_paragraph(ftf, "Example Calculation", font_size=Pt(15), bold=True,
                   color=RGBColor(0x1F, 0x49, 0x7D))
add_body_paragraph(ftf, "Input load: 300 lb  \u00d7  FoS 1.5  =  450 lb design load",
                   font_size=Pt(13), font_name=MONO_FONT)
add_body_paragraph(ftf, "If ultimate strength = 600 lb:  Margin = (600 / 450) - 1 = 0.33 > 0 \u2713",
                   font_size=Pt(13), font_name=MONO_FONT)

scope_box = add_textbox(slide, CONTENT_LEFT, CONTENT_TOP + 2500000,
                        CONTENT_WIDTH, 1000000)
scf = scope_box.text_frame
scf.word_wrap = True

add_body_paragraph(scf, "Applied globally across: actuator mounts, pivot arms, flanges, mounting brackets",
                   font_size=Pt(13), space_before=Pt(0))
add_body_paragraph(scf, "Correction: Previous CDR stated FoS of 1 \u2014 this has been updated to 1.25/1.5 per standard practice.",
                   font_size=Pt(12), italic=True, bold=True, color=RGBColor(0xCC, 0x33, 0x00))

offset += 1

# ============================================================
# 7. FMEA Summary (insert after original slide 35 = Failure Response)
# ============================================================
insert_idx = 35 + offset  # 41
slide = insert_slide_at(prs, insert_idx, LAYOUT)
set_title(slide, "Safety and Failure Behavior: FMEA Summary")

body = clear_body_placeholder(slide)
sp = body._element
sp.getparent().remove(sp)

sub_box = add_textbox(slide, CONTENT_LEFT, CONTENT_TOP - 200000, CONTENT_WIDTH, 300000)
stf = sub_box.text_frame
p = stf.paragraphs[0]
run = p.add_run()
run.text = "Top 10 Highest-RPN Failure Modes (of 45 total)"
run.font.size = Pt(14)
run.font.bold = True
run.font.name = BODY_FONT

tbl = add_table(slide, 11, 5,
                left=CONTENT_LEFT, top=CONTENT_TOP + 100000,
                width=CONTENT_WIDTH, height=Emu(3200000))

fmea_headers = ["ID", "Failure Mode", "RPN", "Risk Level", "Recommended Action"]
for i, h in enumerate(fmea_headers):
    set_cell(tbl, 0, i, h, font_size=TABLE_HEADER_FONT_SIZE, bold=True)
shade_row(tbl, 0, RGBColor(0x1F, 0x49, 0x7D))
for ci in range(5):
    for para in tbl.cell(0, ci).text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

fmea_data = [
    ["E-001", "Water/mud ingress to electronics", "336", "HIGH", "IP67 enclosure, conformal coat, sealed connectors"],
    ["E-002", "Excessive vibration damage", "336", "HIGH", "Vibration-dampening mounts, locking connectors"],
    ["E-003", "Heat from engine bay", "180", "MED", "Heat shield, high-temp wire routing"],
    ["E-004", "EMI from alternator/ignition", "180", "MED", "Shielded cables, ferrite cores, bypass caps"],
    ["R-001", "Relay K1 contacts welded (extend)", "162", "MED", "Size relays for stall current, add feedback monitoring"],
    ["R-002", "Relay K2 contacts welded (retract)", "162", "MED", "Size relays for stall current, add feedback monitoring"],
    ["PW-004", "12V brownout during stall", "160", "MED", "100\u00b5F bulk capacitor, graceful startup sequence"],
    ["H-002", "Intermittent hall sensor connection", "150", "MED", "Deutsch DT connectors, median filter on RPM"],
    ["PW-006", "Connector corrosion", "150", "MED", "Sealed connectors, dielectric grease"],
    ["INT-003", "Actuator-to-sheave linkage failure", "144", "MED", "Linkage FoS > 3, locking hardware"],
]

for r, row_data in enumerate(fmea_data):
    for c, val in enumerate(row_data):
        set_cell(tbl, r + 1, c, val, font_size=Pt(10))
    rpn = int(row_data[2])
    if rpn >= 300:
        shade_row(tbl, r + 1, RGBColor(0xFF, 0xCC, 0xCC))
    elif rpn >= 160:
        shade_row(tbl, r + 1, RGBColor(0xFF, 0xE6, 0xCC))
    else:
        shade_row(tbl, r + 1, RGBColor(0xFF, 0xFF, 0xCC))

fmea_col_widths = [0.08, 0.28, 0.07, 0.10, 0.47]
for i, pct in enumerate(fmea_col_widths):
    tbl.columns[i].width = int(CONTENT_WIDTH * pct)

note_box = add_textbox(slide, CONTENT_LEFT, CONTENT_TOP + 3400000, CONTENT_WIDTH, 300000)
ntf = note_box.text_frame
add_body_paragraph(ntf, "Full 45-item FMEA available in supplementary documentation.",
                   font_size=NOTE_FONT_SIZE, italic=True, color=RGBColor(0x66, 0x66, 0x66),
                   space_before=Pt(0))

offset += 1

# ============================================================
# 8. Budget (insert before Appendix = before original slide 44)
# ============================================================
insert_idx = 43 + offset  # 50
slide = insert_slide_at(prs, insert_idx, LAYOUT)
set_title(slide, "Project Budget")

body = clear_body_placeholder(slide)
sp = body._element
sp.getparent().remove(sp)

tbl = add_table(slide, 13, 4,
                left=CONTENT_LEFT, top=CONTENT_TOP,
                width=int(CONTENT_WIDTH * 0.75), height=Emu(3400000))

budget_headers = ["Item", "Qty", "Unit Cost", "Total"]
for i, h in enumerate(budget_headers):
    set_cell(tbl, 0, i, h, font_size=TABLE_HEADER_FONT_SIZE, bold=True,
             alignment=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)
shade_row(tbl, 0, RGBColor(0x1F, 0x49, 0x7D))
for ci in range(4):
    for para in tbl.cell(0, ci).text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

budget_data = [
    ["Teensy 4.1", "1", "$30", "$30"],
    ["Raspberry Pi 4B", "1", "$55", "$55"],
    ["PQ12-100-6-R Actuator", "1", "$75", "$75"],
    ["Dual Relay Module", "1", "$8", "$8"],
    ["Buck Converter 12V\u21925V", "1", "$12", "$12"],
    ["Hall Effect Sensor", "2", "$3", "$6"],
    ["Potentiometer 10k\u03a9", "1", "$2", "$2"],
    ["LoRa RFM95W Module", "2", "$15", "$30"],
    ["Wiring / Connectors / Fuses", "1", "$40", "$40"],
    ["Enclosure (IP67)", "1", "$25", "$25"],
    ["Misc (resistors, caps, heatshrink)", "1", "$15", "$15"],
    ["TOTAL", "", "", "$298"],
]

for r, row_data in enumerate(budget_data):
    for c, val in enumerate(row_data):
        is_total = (r == len(budget_data) - 1)
        align = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
        set_cell(tbl, r + 1, c, val, bold=is_total, alignment=align,
                 font_size=TABLE_FONT_SIZE if not is_total else TABLE_HEADER_FONT_SIZE)
    if r % 2 == 0:
        shade_row(tbl, r + 1, RGBColor(0xE8, 0xEE, 0xF4))

shade_row(tbl, 12, RGBColor(0x1F, 0x49, 0x7D))
for ci in range(4):
    for para in tbl.cell(12, ci).text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

budget_col_widths = [0.50, 0.12, 0.19, 0.19]
for i, pct in enumerate(budget_col_widths):
    tbl.columns[i].width = int(CONTENT_WIDTH * 0.75 * pct)

bnote_box = add_textbox(slide, CONTENT_LEFT + int(CONTENT_WIDTH * 0.77),
                        CONTENT_TOP, int(CONTENT_WIDTH * 0.23), 1500000)
bntf = bnote_box.text_frame
bntf.word_wrap = True

p = bntf.paragraphs[0]
run = p.add_run()
run.text = "Budget Notes"
run.font.size = Pt(14)
run.font.bold = True
run.font.name = BODY_FONT

add_body_paragraph(bntf, "Total: $298", font_size=Pt(16), bold=True,
                   color=RGBColor(0x1F, 0x49, 0x7D))
add_body_paragraph(bntf, "All components are COTS (commercial off-the-shelf).", font_size=Pt(11))
add_body_paragraph(bntf, "Prices reflect typical online retail (2025-2026).", font_size=Pt(11))
add_body_paragraph(bntf, "Does not include team labor or 3D-printed enclosure material.", font_size=Pt(11))

offset += 1

# ============================================================
# SAVE
# ============================================================
prs.save(OUTPUT_PATH)

print(f"Done. Total slides now: {len(prs.slides)}")
print(f"Added {offset} new slides.")
